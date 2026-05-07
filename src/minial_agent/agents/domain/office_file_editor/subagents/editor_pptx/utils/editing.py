from pathlib import Path
from typing import Any

from pptx import Presentation


def apply_pptx_edit(
    *,
    path: Path,
    operation: str,
    slots: dict[str, str],
    source_filename: str,
) -> list[dict[str, Any]]:
    presentation = Presentation(path)
    page = int(slots.get("PAGE", str(len(presentation.slides) + 1)))
    changed_count = 0
    if operation == "add_slide":
        layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(layout)
        title = slide.shapes.title
        if title is not None:
            title.text = slots.get("TITLE", "")
            changed_count += 1
        body = slots.get("TEXT")
        if body and len(slide.placeholders) > 1:
            slide.placeholders[1].text = body
            changed_count += 1
        new_value = slots.get("TITLE") or body or ""
    else:
        if page < 1 or page > len(presentation.slides):
            raise ValueError(f"PPTX page is out of range: {page}")
        slide = presentation.slides[page - 1]
        new_value = slots.get("TEXT") or slots.get("TITLE")
        if new_value is None:
            raise ValueError("PPTX edit requires TEXT or TITLE.")
        if operation == "replace_slide_title":
            if slide.shapes.title is None:
                raise ValueError(f"PPTX page has no title placeholder: {page}")
            slide.shapes.title.text = new_value
            changed_count = 1
        elif operation == "replace_slide_text":
            old_text = slots.get("OLD_TEXT")
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                if old_text:
                    if old_text in shape.text:
                        shape.text = shape.text.replace(old_text, new_value)
                        changed_count += 1
                elif shape.text:
                    shape.text = new_value
                    changed_count = 1
                    break
        else:
            raise ValueError(f"Unsupported PPTX edit operation: {operation}")
    if changed_count == 0:
        raise ValueError("PPTX edit did not change the presentation.")
    presentation.save(path)
    return [
        {
            "source_file": source_filename,
            "page_number": page,
            "operation": operation,
            "new_value": new_value,
            "changed_count": changed_count,
        }
    ]
