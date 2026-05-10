from pathlib import Path

from fastapi import FastAPI
from fastapi.testclient import TestClient
from openpyxl import Workbook
from pptx import Presentation

from minial_agent.api.fs.router import router as fs_router
from minial_agent.integrations.upload import ensure_upload_workspace


def test_preview_metadata_for_supported_file_types(tmp_path, monkeypatch) -> None:
    monkeypatch.setenv("AGENT_RUNTIME_ROOT_DIR", str(tmp_path))
    monkeypatch.setattr(
        "minial_agent.integrations.fs.preview.convert_to_pdf",
        _fake_convert_to_pdf,
    )
    monkeypatch.setattr(
        "minial_agent.integrations.fs.preview.build_pptx_preview",
        lambda _path: {"slide_count": 1, "slides": []},
    )
    workspace = ensure_upload_workspace(tmp_path / "user")
    (workspace.files_dir / "sample.pdf").write_bytes(b"%PDF-1.7\n")
    (workspace.files_dir / "report.docx").write_bytes(b"docx")
    (workspace.files_dir / "slides.pptx").write_bytes(b"pptx")
    (workspace.files_dir / "form.hwpx").write_bytes(b"hwpx")
    (workspace.files_dir / "README.md").write_text("# Title\n", encoding="utf-8")
    (workspace.files_dir / "notes.markdown").write_text("Notes", encoding="utf-8")
    (workspace.files_dir / "plain.txt").write_text("plain", encoding="utf-8")
    (workspace.files_dir / "script.py").write_text("print('ok')\n", encoding="utf-8")
    (workspace.files_dir / "app.js").write_text("console.log('ok')\n", encoding="utf-8")
    (workspace.files_dir / "app.ts").write_text("const ok: boolean = true\n", encoding="utf-8")
    (workspace.files_dir / "query.sql").write_text("select 1;\n", encoding="utf-8")
    (workspace.files_dir / "index.html").write_text("<main></main>\n", encoding="utf-8")
    (workspace.files_dir / "style.css").write_text("body { margin: 0; }\n", encoding="utf-8")
    (workspace.files_dir / "Main.java").write_text("class Main {}\n", encoding="utf-8")
    (workspace.files_dir / "main.go").write_text("package main\n", encoding="utf-8")
    (workspace.files_dir / "run.sh").write_text("echo ok\n", encoding="utf-8")
    (workspace.files_dir / "data.json").write_text('{"ok": true}\n', encoding="utf-8")
    _write_workbook(workspace.files_dir / "book.xlsx")

    client = _client()

    expected = {
        "sample.pdf": "pdf",
        "report.docx": "office_pdf",
        "slides.pptx": "office_pdf",
        "form.hwpx": "hwpx",
        "book.xlsx": "xlsx_grid",
        "README.md": "markdown",
        "notes.markdown": "markdown",
        "plain.txt": "text",
        "script.py": "code",
        "app.js": "code",
        "app.ts": "code",
        "query.sql": "code",
        "index.html": "code",
        "style.css": "code",
        "Main.java": "code",
        "main.go": "code",
        "run.sh": "code",
        "data.json": "code",
    }
    for filename, preview_type in expected.items():
        response = client.get(
            "/api/fs/preview",
            params={"user_id": "user", "uuid": "session", "path": f"files/{filename}"},
        )

        assert response.status_code == 200
        body = response.json()
        assert body["path"] == f"files/{filename}"
        assert body["filename"] == filename
        assert body["preview_type"] == preview_type
        if filename == "slides.pptx":
            assert body["presentation"]["slide_count"] == 1


def test_preview_rejects_internal_paths(tmp_path, monkeypatch) -> None:
    monkeypatch.setenv("AGENT_RUNTIME_ROOT_DIR", str(tmp_path))
    ensure_upload_workspace(tmp_path / "user")

    response = _client().get(
        "/api/fs/preview",
        params={"user_id": "user", "uuid": "session", "path": "../.registry/files.json"},
    )

    assert response.status_code == 400


def test_preview_source_serves_public_pdf(tmp_path, monkeypatch) -> None:
    monkeypatch.setenv("AGENT_RUNTIME_ROOT_DIR", str(tmp_path))
    workspace = ensure_upload_workspace(tmp_path / "user")
    (workspace.files_dir / "sample.pdf").write_bytes(b"%PDF-1.7\n")

    response = _client().get(
        "/api/fs/preview/source",
        params={"user_id": "user", "uuid": "session", "path": "files/sample.pdf"},
    )

    assert response.status_code == 200
    assert response.headers["content-type"] == "application/pdf"
    assert response.content == b"%PDF-1.7\n"


def test_preview_source_serves_text_and_markdown(tmp_path, monkeypatch) -> None:
    monkeypatch.setenv("AGENT_RUNTIME_ROOT_DIR", str(tmp_path))
    workspace = ensure_upload_workspace(tmp_path / "user")
    (workspace.files_dir / "README.md").write_text("# Title\n", encoding="utf-8")
    (workspace.files_dir / "plain.txt").write_text("plain", encoding="utf-8")

    client = _client()
    markdown_response = client.get(
        "/api/fs/preview/source",
        params={"user_id": "user", "uuid": "session", "path": "files/README.md"},
    )
    text_response = client.get(
        "/api/fs/preview/source",
        params={"user_id": "user", "uuid": "session", "path": "files/plain.txt"},
    )

    assert markdown_response.status_code == 200
    assert markdown_response.headers["content-type"] == "text/markdown; charset=utf-8"
    assert markdown_response.text == "# Title\n"
    assert text_response.status_code == 200
    assert text_response.headers["content-type"] == "text/plain; charset=utf-8"
    assert text_response.text == "plain"


def test_preview_source_serves_code_as_plain_text(tmp_path, monkeypatch) -> None:
    monkeypatch.setenv("AGENT_RUNTIME_ROOT_DIR", str(tmp_path))
    workspace = ensure_upload_workspace(tmp_path / "user")
    (workspace.files_dir / "script.py").write_text("print('ok')\n", encoding="utf-8")

    response = _client().get(
        "/api/fs/preview/source",
        params={"user_id": "user", "uuid": "session", "path": "files/script.py"},
    )

    assert response.status_code == 200
    assert response.headers["content-type"] == "text/plain; charset=utf-8"
    assert response.text == "print('ok')\n"


def test_xlsx_preview_returns_grid_data_and_sheet_order(tmp_path, monkeypatch) -> None:
    monkeypatch.setenv("AGENT_RUNTIME_ROOT_DIR", str(tmp_path))
    workspace = ensure_upload_workspace(tmp_path / "user")
    _write_workbook(workspace.files_dir / "book.xlsx")

    response = _client().get(
        "/api/fs/preview",
        params={"user_id": "user", "uuid": "session", "path": "files/book.xlsx"},
    )

    assert response.status_code == 200
    workbook = response.json()["workbook"]
    assert workbook["sheet_count"] == 2
    assert [sheet["name"] for sheet in workbook["sheets"]] == ["Summary", "Details"]
    assert workbook["sheets"][0]["merged_ranges"] == ["A1:B1"]
    assert workbook["sheets"][0]["cells"][0]["address"] == "A1"
    assert workbook["sheets"][0]["cells"][0]["value"] == "Revenue"
    assert workbook["sheets"][1]["cells"][0]["value"] == "Region"


def test_pptx_operations_change_editable_shape(tmp_path, monkeypatch) -> None:
    monkeypatch.setenv("AGENT_RUNTIME_ROOT_DIR", str(tmp_path))
    monkeypatch.setattr(
        "minial_agent.integrations.fs.preview.convert_to_pdf",
        _fake_convert_to_pdf,
    )
    workspace = ensure_upload_workspace(tmp_path / "user")
    path = workspace.files_dir / "slides.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Old title"
    title_shape_id = slide.shapes.title.shape_id
    presentation.save(path)

    deck_response = _client().get(
        "/api/fs/pptx/deck",
        params={"user_id": "user", "uuid": "session", "path": "files/slides.pptx"},
    )
    assert deck_response.status_code == 200
    deck = deck_response.json()["deck"]
    element_id = f"shape-{title_shape_id}"

    response = _client().post(
        "/api/fs/pptx/operations",
        json={
            "user_id": "user",
            "uuid": "session",
            "path": "files/slides.pptx",
            "origin": "user",
            "expected_revision": deck["revision"],
            "operations": [
                {
                    "type": "updateText",
                    "slideId": "slide-1",
                    "elementId": element_id,
                    "content": "New title",
                },
                {
                    "type": "moveElement",
                    "slideId": "slide-1",
                    "elementId": element_id,
                    "x": 457200,
                    "y": 365760,
                },
                {
                    "type": "resizeElement",
                    "slideId": "slide-1",
                    "elementId": element_id,
                    "width": 2743200,
                    "height": 914400,
                },
            ],
        },
    )

    assert response.status_code == 200
    assert response.json()["revision"] == 1
    edited = Presentation(path)
    edited_title = edited.slides[0].shapes.title
    assert edited_title.text == "New title"
    assert int(edited_title.left) == 457200
    assert int(edited_title.top) == 365760
    assert int(edited_title.width) == 2743200
    assert int(edited_title.height) == 914400


def test_pptx_search_uses_fts_index(tmp_path, monkeypatch) -> None:
    monkeypatch.setenv("AGENT_RUNTIME_ROOT_DIR", str(tmp_path))
    workspace = ensure_upload_workspace(tmp_path / "user")
    path = workspace.files_dir / "slides.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Quarterly review"
    slide.placeholders[1].text = "Revenue grew 12 percent"
    presentation.save(path)

    response = _client().get(
        "/api/fs/pptx/search",
        params={
            "user_id": "user",
            "uuid": "session",
            "path": "files/slides.pptx",
            "q": "Revenue",
        },
    )

    assert response.status_code == 200
    assert response.json()["matches"][0]["slideId"] == "slide-1"


def test_pptx_ai_operation_respects_manual_override(tmp_path, monkeypatch) -> None:
    monkeypatch.setenv("AGENT_RUNTIME_ROOT_DIR", str(tmp_path))
    workspace = ensure_upload_workspace(tmp_path / "user")
    path = workspace.files_dir / "slides.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Old title"
    title_shape_id = slide.shapes.title.shape_id
    presentation.save(path)

    deck = _client().get(
        "/api/fs/pptx/deck",
        params={"user_id": "user", "uuid": "session", "path": "files/slides.pptx"},
    ).json()["deck"]
    element_id = f"shape-{title_shape_id}"
    user_response = _client().post(
        "/api/fs/pptx/operations",
        json={
            "user_id": "user",
            "uuid": "session",
            "path": "files/slides.pptx",
            "origin": "user",
            "expected_revision": deck["revision"],
            "operations": [
                {
                    "type": "updateText",
                    "slideId": "slide-1",
                    "elementId": element_id,
                    "content": "Manual title",
                },
            ],
        },
    )
    assert user_response.status_code == 200

    ai_response = _client().post(
        "/api/fs/pptx/operations",
        json={
            "user_id": "user",
            "uuid": "session",
            "path": "files/slides.pptx",
            "origin": "ai",
            "expected_revision": user_response.json()["revision"],
            "operations": [
                {
                    "type": "updateText",
                    "slideId": "slide-1",
                    "elementId": element_id,
                    "content": "AI title",
                },
            ],
        },
    )

    assert ai_response.status_code == 200
    body = ai_response.json()
    assert body["changed_slide_ids"] == []
    assert body["rejected_operations"]
    assert Presentation(path).slides[0].shapes.title.text == "Manual title"


def _fake_convert_to_pdf(source_path: Path, output_dir: Path, target_pdf: Path) -> None:
    output_dir.mkdir(parents=True, exist_ok=True)
    target_pdf.parent.mkdir(parents=True, exist_ok=True)
    target_pdf.write_bytes(b"%PDF-1.7\nconverted " + source_path.name.encode())


def _write_workbook(path: Path) -> None:
    workbook = Workbook()
    summary = workbook.active
    summary.title = "Summary"
    summary.merge_cells("A1:B1")
    summary["A1"] = "Revenue"
    summary["A2"] = "Q1"
    summary["B2"] = 10
    details = workbook.create_sheet("Details")
    details["A1"] = "Region"
    details["B1"] = "Amount"
    details["A2"] = "KR"
    details["B2"] = 7
    workbook.save(path)
    workbook.close()


def _client() -> TestClient:
    app = FastAPI()
    app.include_router(fs_router)
    return TestClient(app)
