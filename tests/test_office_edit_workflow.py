import zipfile
import json
from pathlib import Path

import pytest
from docx import Document
from pptx import Presentation

from minial_agent.agents.domain.office_file_editor.subagents.editor_docx.utils.editing import (
    apply_docx_edit,
)
from minial_agent.agents.domain.office_file_editor.subagents.editor_docx.workflow.edit.nodes import (
    build_docx_edit_spec,
)
from minial_agent.agents.domain.office_file_editor.subagents.editor_hwpx.utils.editing import (
    apply_hwpx_edit,
)
from minial_agent.agents.domain.office_file_editor.subagents.editor_pptx.utils.editing import (
    apply_pptx_edit,
)
from minial_agent.agents.domain.office_file_editor.subagents.editor_pptx.workflow.edit.nodes import (
    build_pptx_edit_spec,
)
from minial_agent.integrations.pptx.preview import build_pptx_preview
from minial_agent.integrations.upload import ensure_upload_workspace
from minial_agent.integrations.upload.registry import UploadRegistry
from minial_agent.integrations.upload.resolver import resolve_upload_artifact
from minial_agent.agents.domain.office_file_editor.utils.edit_protocol import (
    parse_slots,
    select_operation,
)


def test_edit_protocol_rejects_malformed_outputs() -> None:
    with pytest.raises(ValueError, match="one line"):
        select_operation(
            instruction="edit",
            prompt="{instruction}",
            selector=lambda _instruction: "replace_text\nextra",
        )

    with pytest.raises(ValueError, match="JSON or markdown"):
        parse_slots('{"OLD_TEXT": "a"}')

    with pytest.raises(ValueError, match="KEY=VALUE"):
        parse_slots("OLD_TEXT=a; bad_segment")

    with pytest.raises(ValueError, match="non-empty"):
        parse_slots("OLD_TEXT=a; NEW_TEXT=")


def test_docx_edit_operations_change_file(tmp_path) -> None:
    path = tmp_path / "report.docx"
    document = Document()
    document.add_paragraph("FY2024 revenue")
    table = document.add_table(rows=1, cols=1)
    table.cell(0, 0).text = "old cell"
    document.save(path)

    replace_result = apply_docx_edit(
        path=path,
        operation="replace_text",
        slots={"OLD_TEXT": "FY2024", "NEW_TEXT": "FY2025"},
        source_filename="report.docx",
    )
    add_result = apply_docx_edit(
        path=path,
        operation="add_paragraph",
        slots={"TEXT": "new paragraph"},
        source_filename="report.docx",
    )
    table_result = apply_docx_edit(
        path=path,
        operation="update_table_cell",
        slots={"TABLE": "1", "ROW": "1", "COLUMN": "1", "TEXT": "new cell"},
        source_filename="report.docx",
    )

    edited = Document(path)
    assert edited.paragraphs[0].text == "FY2025 revenue"
    assert edited.paragraphs[-1].text == "new paragraph"
    assert edited.tables[0].cell(0, 0).text == "new cell"
    assert replace_result[0]["changed_count"] == 1
    assert add_result[0]["changed_count"] == 1
    assert table_result[0]["changed_count"] == 1


def test_docx_edit_rejects_missing_slot() -> None:
    with pytest.raises(ValueError, match="NEW_TEXT"):
        build_docx_edit_spec(
            instruction="replace",
            operation_selector=lambda _instruction: "replace_text",
            slot_filler=lambda _operation, _instruction: "OLD_TEXT=FY2024",
        )


def test_pptx_edit_operations_change_file(tmp_path) -> None:
    path = tmp_path / "deck.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Old title"
    slide.placeholders[1].text = "Old body"
    presentation.save(path)

    apply_pptx_edit(
        path=path,
        operation="replace_slide_title",
        slots={"PAGE": "1", "TITLE": "New title"},
        source_filename="deck.pptx",
    )
    apply_pptx_edit(
        path=path,
        operation="replace_slide_text",
        slots={"PAGE": "1", "OLD_TEXT": "Old body", "TEXT": "New body"},
        source_filename="deck.pptx",
    )
    add_result = apply_pptx_edit(
        path=path,
        operation="add_slide",
        slots={"TITLE": "Added slide", "TEXT": "Added body"},
        source_filename="deck.pptx",
    )

    edited = Presentation(path)
    assert edited.slides[0].shapes.title.text == "New title"
    assert edited.slides[0].placeholders[1].text == "New body"
    assert edited.slides[1].shapes.title.text == "Added slide"
    assert add_result[0]["page_number"] == 2


def test_pptx_preview_extracts_slide_outline(tmp_path) -> None:
    path = tmp_path / "deck.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Quarterly review"
    slide.placeholders[1].text = "Revenue grew 12%"
    presentation.save(path)

    preview = build_pptx_preview(path)

    assert preview["canvas"]["width"] > 0
    assert preview["canvas"]["height"] > 0
    assert preview["revision"] == 0
    assert len(preview["slides"]) == 1
    assert preview["slides"][0]["title"] == "Quarterly review"
    assert preview["slides"][0]["elements"][0]["content"] == "Quarterly review"
    assert preview["slides"][0]["elements"][0]["width"] > 0


def test_pptx_ai_edit_spec_uses_operation_json(tmp_path) -> None:
    workspace = ensure_upload_workspace(tmp_path)
    path = workspace.files_dir / "deck.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Old title"
    title_shape_id = slide.shapes.title.shape_id
    presentation.save(path)
    converted_dir = workspace.converted_dir / "file_001"
    converted_dir.mkdir(parents=True)
    registry = UploadRegistry(workspace.registry_path)
    registry.add_uploaded(
        file_id="file_001",
        visible_path=path,
        visible_name=path.name,
        file_type="pptx",
        converted_dir=converted_dir,
    )
    registry.update_status("file_001", status="converted")
    (converted_dir / "manifest.json").write_text(
        json.dumps(
            {
                "file_id": "file_001",
                "source_filename": path.name,
                "source_path": str(path),
                "file_type": "pptx",
                "converted_dir": str(converted_dir),
                "pages": [],
                "status": "converted",
            }
        ),
        encoding="utf-8",
    )
    artifact = resolve_upload_artifact(workspace=workspace, file_ref="deck.pptx")

    spec = build_pptx_edit_spec(
        workspace=workspace,
        artifact=artifact,
        instruction="change title",
        operation_generator=lambda _instruction, _deck: (
            '[{"type":"updateText","slideId":"slide-1",'
            f'"elementId":"shape-{title_shape_id}","content":"New title"}}]'
        ),
    )

    assert spec["operations"][0]["type"] == "updateText"
    assert spec["operations"][0]["elementId"] == f"shape-{title_shape_id}"


def test_hwpx_edit_operations_change_zip_xml(tmp_path) -> None:
    path = tmp_path / "proposal.hwpx"
    _write_hwpx(path, "<hp:body><hp:p>OLD</hp:p></hp:body>")

    apply_hwpx_edit(
        path=path,
        operation="replace_text",
        slots={"OLD_TEXT": "OLD", "NEW_TEXT": "NEW"},
        source_filename="proposal.hwpx",
    )
    apply_hwpx_edit(
        path=path,
        operation="add_paragraph",
        slots={"TEXT": "added text"},
        source_filename="proposal.hwpx",
    )

    with zipfile.ZipFile(path) as package:
        xml = package.read("Contents/section0.xml").decode("utf-8")
    assert "NEW" in xml
    assert "added text" in xml


def _write_hwpx(path: Path, xml: str) -> None:
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as package:
        package.writestr("Contents/section0.xml", xml)
